import os
import io
import json
import pickle
import sqlite3
import time
from io import BytesIO
from datetime import datetime, timedelta
from urllib.parse import quote_plus

import streamlit as st
import pdfplumber
from gtts import gTTS
import httpx, certifi
from dotenv import load_dotenv
import numpy as np

try:
    import plotly.express as px
    import plotly.graph_objects as go
    PLOTLY_AVAILABLE = True
except:
    PLOTLY_AVAILABLE = False

try:
    import pypdfium2 as pdfium
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False

st.set_page_config(page_title="AI Mentor Study Assistant Pro", layout="wide")
load_dotenv()
os.environ.setdefault("SSL_CERT_FILE", certifi.where())

try:
    from openai import OpenAI
    _openai_key = os.environ.get("OPENAI_API_KEY") or (st.secrets.get("OPENAI_API_KEY") if hasattr(st, "secrets") else None)
    client = OpenAI(api_key=_openai_key) if _openai_key else None
except Exception:
    client = None

# ============== DATABASE SETUP ==============
DB_FILE = "study_sessions.db"

def init_db():
    """Initialize SQLite database for tracking progress"""
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS quiz_scores
                 (id INTEGER PRIMARY KEY, session_name TEXT, topic TEXT, score REAL, total INTEGER, timestamp TEXT)''')
    c.execute('''CREATE TABLE IF NOT EXISTS study_sessions
                 (id INTEGER PRIMARY KEY, name TEXT, topics TEXT, summary TEXT, created_at TEXT, updated_at TEXT)''')
    c.execute('''CREATE TABLE IF NOT EXISTS achievements
                 (id INTEGER PRIMARY KEY, session_id INTEGER, badge_name TEXT, achieved_at TEXT)''')
    conn.commit()
    conn.close()

def add_quiz_score(session_name, topic, score, total):
    """Save quiz score to database"""
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("INSERT INTO quiz_scores (session_name, topic, score, total, timestamp) VALUES (?, ?, ?, ?, ?)",
              (session_name, topic, score, total, datetime.now().isoformat()))
    conn.commit()
    conn.close()

def get_quiz_history(session_name=None):
    """Retrieve quiz history"""
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    if session_name:
        c.execute("SELECT * FROM quiz_scores WHERE session_name = ? ORDER BY timestamp DESC", (session_name,))
    else:
        c.execute("SELECT * FROM quiz_scores ORDER BY timestamp DESC")
    results = c.fetchall()
    conn.close()
    return results

# ============== STYLING ==============
st.markdown("""
<style>
.status-badge {
    display: inline-block;
    padding: 8px 12px;
    border-radius: 20px;
    font-weight: bold;
    font-size: 0.9rem;
}
.status-badge.done {
    background-color: #10b981;
    color: white;
}
.status-badge.pending {
    background-color: #fbbf24;
    color: #1a1a1a;
}
.flashcard {
    border: 2px solid #6366f1;
    border-radius: 12px;
    padding: 20px;
    text-align: center;
    min-height: 200px;
    display: flex;
    align-items: center;
    justify-content: center;
    cursor: pointer;
    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    color: white;
    font-size: 18px;
    font-weight: bold;
}
.achievement {
    display: inline-block;
    padding: 12px;
    margin: 5px;
    border-radius: 8px;
    background-color: #fbbf24;
    color: #1a1a1a;
    font-weight: bold;
}
</style>
""", unsafe_allow_html=True)

st.markdown("""<div style="text-align:center; padding: 32px 0;">
<h1 style="margin:0;">🎓 <span style="color:#6366f1;">AI Mentor Study Assistant Pro</span></h1>
<p style="color:#666;">Smart Learning • Flashcards • Progress Tracking • Q&A • Exams • Achievements</p>
</div>""", unsafe_allow_html=True)

st.markdown("---")

# ============== HELPER FUNCTIONS ==============
def fetch_videos_for_topics(topics):
    """Generate video resource links"""
    results = {}
    for t in topics:
        q = quote_plus(f"{t} tutorial")
        url = f"https://www.youtube.com/results?search_query={q}"
        results[t] = [{"title": f"📺 {t}", "url": url, "note": "YouTube search"}]
    return results

def generate_lab_content(text: str, topic: str) -> dict:
    """Generate lab-related content: project ideas, hands-on assignments, lab tips, and real-world applications."""
    if not text or not topic:
        return {"projects": [], "assignments": [], "lab_tips": [], "real_world": ""}
    
    if not client:
        return {"projects": [], "assignments": [], "lab_tips": [], "real_world": ""}
    
    prompt = f'''Based on this topic: "{topic}"

From the provided text context, suggest:
1. 3 hands-on project ideas (real-world applications) that demonstrate this topic in action
2. 5 practical lab assignments or experiments students can do
3. 4 lab tips for successful hands-on work
4. Brief explanation of how this topic applies to real-world projects

Format your response as JSON with keys: "projects" (array), "assignments" (array), "lab_tips" (array), "real_world" (string).
Return ONLY JSON.

Text context: {truncate_chars(text, 3000)}'''
    
    response = openai_complete(prompt, temperature=0.3, max_tokens=1000)
    try:
        if response:
            lab_data = json.loads(response.replace("```json","").replace("```","").strip())
            return lab_data
    except Exception as e:
        st.error(f"Lab content generation failed: {e}")
    
    return {"projects": [], "assignments": [], "lab_tips": [], "real_world": ""}

def truncate_chars(s: str, max_chars: int = 12000) -> str:
    """Truncate text to max length"""
    return s[:max_chars] if s else ""

@st.cache_data(show_spinner=False)
def extract_pdf_text(uploaded_file, use_ocr=False) -> str:
    """Extract text from PDF"""
    if not uploaded_file:
        return ""
    raw = uploaded_file.getvalue()
    text = ""
    try:
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t:
                    text += t + "\n"
    except Exception:
        pass
    return text.strip()


# ---------- alternative provider helpers ----------

def ollama_generate(prompt: str, model: str = "llama3") -> str:
    """Use a local Ollama server (must be running with `ollama serve`)."""
    try:
        r = httpx.post(
            "http://localhost:11434/api/generate",
            json={"model": model, "prompt": prompt, "stream": False},
            timeout=120,
        )
        r.raise_for_status()
        return r.json().get("response", "").strip()
    except Exception as e:
        st.error(f"Ollama error: {e}")
        return None


def hf_generate(prompt: str) -> str:
    """Call HuggingFace inference API. Requires HF_TOKEN in env or secrets."""
    hf_token = os.environ.get("HF_TOKEN") or (
        st.secrets.get("HF_TOKEN") if hasattr(st, "secrets") else None
    )
    if not hf_token:
        st.error("⚠️ HuggingFace token not configured!")
        return None
    model_name = st.session_state.get("hf_model", "google/flan-t5-large")
    url = f"https://api-inference.huggingface.co/models/{model_name}"
    headers = {"Authorization": f"Bearer {hf_token}"}
    try:
        r = httpx.post(url, headers=headers, json={"inputs": prompt}, timeout=120)
        r.raise_for_status()
        data = r.json()
        if isinstance(data, list) and data and "generated_text" in data[0]:
            return data[0]["generated_text"].strip()
        return str(data)
    except Exception as e:
        st.error(f"HuggingFace API error: {e}")
        return None


def openai_complete(prompt: str, temperature=0.25, max_tokens=900) -> str:
    """Unified LLM interface selecting provider via session_state.provider."""
    provider = st.session_state.get("provider", "openai")
    if provider == "openai":
        if not client:
            st.error("⚠️ OpenAI API key not configured!")
            return None
        try:
            resp = client.chat.completions.create(
                model=st.session_state.get("model", "gpt-4o-mini"),
                messages=[{"role": "user", "content": prompt}],
                temperature=temperature,
                max_tokens=max_tokens,
            )
            return (resp.choices[0].message.content or "").strip()
        except Exception as e:
            st.error(f"OpenAI API Error: {str(e)}")
            return None
    elif provider == "ollama":
        return ollama_generate(prompt, model=st.session_state.get("model", "llama3"))
    elif provider == "huggingface":
        return hf_generate(prompt)
    else:
        st.error(f"Unknown provider: {provider}")
        return None


def ai_summarize(text: str) -> str:
    """Generate a summary; chunk text when very large to avoid input limits."""
    if not text:
        return ""
    max_chars = 6000
    if len(text) <= max_chars:
        prompt = f"Summarize these notes in 300 words, with key concepts and one example:\n\n{text}"
        return openai_complete(prompt, temperature=0.25, max_tokens=500) or ""
    # split into character chunks
    chunks = [text[i : i + max_chars] for i in range(0, len(text), max_chars)]
    summaries = []
    for c in chunks:
        s = openai_complete(
            f"Summarize these notes in 300 words, with key concepts and one example:\n\n{c}",
            temperature=0.25,
            max_tokens=500,
        )
        if s:
            summaries.append(s)
    if not summaries:
        return ""
    combined = "\n\n".join(summaries)
    final_prompt = (
        "Combine the following chunk summaries into a single polished summary:\n\n" + combined
    )
    return openai_complete(final_prompt, temperature=0.25, max_tokens=500) or ""


def ai_extract_topics(text: str) -> list:
    """Extract topics from text"""
    if not text:
        return []
    prompt = f'Extract 3-5 study topics from this text as JSON list. Return only: ["topic1","topic2"]\n\nText:\n{truncate_chars(text, 4000)}'
    try:
        response = openai_complete(prompt, temperature=0.2, max_tokens=200)
        if not response:
            return []
        topics = json.loads(response.replace("```json", "").replace("```", "").strip())
        return [str(t).strip()[:50] for t in topics if t][:6]
    except Exception as e:
        st.error(f"Topic extraction failed: {e}")
        return []

def ai_summarize(text: str) -> str:
    """Generate summary"""
    if not text:
        return ""
    prompt = f"Summarize these notes in 300 words, with key concepts and one example:\n\n{truncate_chars(text, 6000)}"
    return openai_complete(prompt, temperature=0.25, max_tokens=500)

# ============== TIER 1: FLASHCARD GENERATOR ==============
def generate_flashcards(text: str, num_cards=10):
    """Generate flashcards from text"""
    if not text:
        return []
    prompt = f'''Create {num_cards} flashcard pairs (front/back) from this text.
Format as JSON array: [{{"front":"Question?","back":"Answer..."}}]
Return ONLY JSON.

Text: {truncate_chars(text, 5000)}'''
    response = openai_complete(prompt, temperature=0.2, max_tokens=1500)
    try:
        if response:
            flashcards = json.loads(response.replace("```json","").replace("```","").strip())
            return flashcards
    except:
        st.error("Failed to generate flashcards")
    return []

def export_anki_format(flashcards):
    """Export flashcards as Anki format (TSV)"""
    if not flashcards:
        return ""
    lines = []
    for fc in flashcards:
        lines.append(f"{fc['front']}\t{fc['back']}")
    return "\n".join(lines)


def export_to_docx(flashcards, title="Flashcards"):
    """Export flashcards to a DOCX file (returns bytes or (None, error_message))."""
    if not flashcards:
        return None, "No flashcards to export"
    try:
        from docx import Document
    except Exception:
        return None, "Missing dependency: install python-docx (pip install python-docx)"

    try:
        doc = Document()
        doc.add_heading(title, level=1)
        for i, fc in enumerate(flashcards, 1):
            doc.add_heading(f"Card {i}", level=2)
            doc.add_paragraph("Front: ")
            doc.add_paragraph(fc.get("front", ""))
            doc.add_paragraph("Back: ")
            doc.add_paragraph(fc.get("back", ""))

        bio = BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio.getvalue(), None
    except Exception as e:
        return None, f"DOCX export failed: {e}"


def export_to_pdf(flashcards, title="Flashcards"):
    """Export flashcards to a simple PDF (returns bytes or (None, error_message))."""
    if not flashcards:
        return None, "No flashcards to export"
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
    except Exception:
        return None, "Missing dependency: install reportlab (pip install reportlab)"

    try:
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        y = height - inch
        c.setFont("Helvetica-Bold", 16)
        c.drawString(inch, y, title)
        y -= 0.4 * inch
        c.setFont("Helvetica", 11)

        for i, fc in enumerate(flashcards, 1):
            front = f"Card {i} - Front: {fc.get('front','')}"
            back = f"Back: {fc.get('back','')}"
            for line in (front + "\n" + back).split("\n"):
                if y < inch:
                    c.showPage()
                    y = height - inch
                    c.setFont("Helvetica", 11)
                # truncate long lines to avoid overflow
                c.drawString(inch, y, line[:200])
                y -= 14
            y -= 10

        c.save()
        buffer.seek(0)
        return buffer.getvalue(), None
    except Exception as e:
        return None, f"PDF export failed: {e}"


def export_study_pack_to_docx(pack: dict, title_prefix="Study Pack"):
    """Export full study pack (topics, summary, flashcards, quiz) to DOCX. Returns (bytes, error)."""
    try:
        from docx import Document
    except Exception:
        return None, "Missing dependency: install python-docx (pip install python-docx)"

    try:
        doc = Document()
        title = pack.get("session_name") or title_prefix
        doc.add_heading(title, level=1)

        topics = pack.get("topics") or []
        if topics:
            doc.add_heading("Topics", level=2)
            for t in topics:
                doc.add_paragraph(str(t), style='List Bullet')

        summary = pack.get("summary") or ""
        if summary:
            doc.add_heading("Summary", level=2)
            doc.add_paragraph(summary)

        flashcards = pack.get("flashcards") or []
        if flashcards:
            doc.add_heading("Flashcards", level=2)
            for i, fc in enumerate(flashcards, 1):
                doc.add_paragraph(f"Card {i}", style='List Number')
                doc.add_paragraph("Front: ")
                doc.add_paragraph(fc.get('front', ''))
                doc.add_paragraph("Back: ")
                doc.add_paragraph(fc.get('back', ''))

        quiz = pack.get("quiz") or []
        if quiz:
            doc.add_heading("Quiz", level=2)
            for i, q in enumerate(quiz, 1):
                doc.add_paragraph(f"Q{i}. {q.get('question','')}")
                choices = q.get('choices') or []
                for c in choices:
                    doc.add_paragraph(f"- {c}")
                ans_idx = q.get('answer_index')
                if ans_idx is not None and isinstance(ans_idx, int) and 0 <= ans_idx < len(choices):
                    doc.add_paragraph(f"Answer: {choices[ans_idx]}")

        bio = BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio.getvalue(), None
    except Exception as e:
        return None, f"DOCX export failed: {e}"


def export_study_pack_to_pdf(pack: dict, title_prefix="Study Pack"):
    """Export full study pack to PDF. Returns (bytes, error)."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
    except Exception:
        return None, "Missing dependency: install reportlab (pip install reportlab)"

    try:
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        x_margin = inch
        y = height - inch

        title = pack.get("session_name") or title_prefix
        c.setFont("Helvetica-Bold", 16)
        c.drawString(x_margin, y, title)
        y -= 0.4 * inch

        topics = pack.get("topics") or []
        if topics:
            c.setFont("Helvetica-Bold", 12)
            c.drawString(x_margin, y, "Topics:")
            y -= 0.25 * inch
            c.setFont("Helvetica", 10)
            for t in topics:
                if y < inch:
                    c.showPage()
                    y = height - inch
                c.drawString(x_margin + 10, y, f"- {str(t)[:120]}")
                y -= 14
            y -= 8

        summary = pack.get("summary") or ""
        if summary:
            if y < inch:
                c.showPage(); y = height - inch
            c.setFont("Helvetica-Bold", 12)
            c.drawString(x_margin, y, "Summary:")
            y -= 0.25 * inch
            c.setFont("Helvetica", 10)
            for line in summary.split("\n"):
                for chunk in [line[i:i+100] for i in range(0, len(line), 100)]:
                    if y < inch:
                        c.showPage(); y = height - inch
                    c.drawString(x_margin, y, chunk)
                    y -= 12
            y -= 8

        flashcards = pack.get("flashcards") or []
        if flashcards:
            if y < inch:
                c.showPage(); y = height - inch
            c.setFont("Helvetica-Bold", 12)
            c.drawString(x_margin, y, "Flashcards:")
            y -= 0.25 * inch
            c.setFont("Helvetica", 10)
            for i, fc in enumerate(flashcards, 1):
                front = f"Card {i} Front: {fc.get('front','')[:200]}"
                back = f"Back: {fc.get('back','')[:200]}"
                for line in (front + "\n" + back).split("\n"):
                    if y < inch:
                        c.showPage(); y = height - inch
                    c.drawString(x_margin + 10, y, line)
                    y -= 12
                y -= 8

        quiz = pack.get("quiz") or []
        if quiz:
            if y < inch:
                c.showPage(); y = height - inch
            c.setFont("Helvetica-Bold", 12)
            c.drawString(x_margin, y, "Quiz:")
            y -= 0.25 * inch
            c.setFont("Helvetica", 10)
            for i, q in enumerate(quiz, 1):
                question = f"Q{i}. {q.get('question','')[:200]}"
                if y < inch:
                    c.showPage(); y = height - inch
                c.drawString(x_margin + 10, y, question)
                y -= 14
                choices = q.get('choices') or []
                for cchoice in choices:
                    if y < inch:
                        c.showPage(); y = height - inch
                    c.drawString(x_margin + 20, y, f"- {cchoice[:160]}")
                    y -= 12
                ans_idx = q.get('answer_index')
                if ans_idx is not None and isinstance(ans_idx, int) and 0 <= ans_idx < len(choices):
                    if y < inch:
                        c.showPage(); y = height - inch
                    c.drawString(x_margin + 20, y, f"Answer: {choices[ans_idx]}")
                    y -= 14
                y -= 8

        c.save()
        buffer.seek(0)
        return buffer.getvalue(), None
    except Exception as e:
        return None, f"PDF export failed: {e}"


def export_flashcards_to_pptx(flashcards, title="Flashcards"):
    """Export flashcards to a PPTX presentation. Returns (bytes, error)."""
    if not flashcards:
        return None, "No flashcards to export"
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return None, "Missing dependency: install python-pptx (pip install python-pptx)"

    try:
        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title

        for i, fc in enumerate(flashcards, 1):
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            # title
            try:
                s.shapes.title.text = f"Card {i}: {fc.get('front','')[:80]}"
            except Exception:
                pass
            # body
            try:
                body = s.shapes.placeholders[1].text_frame
                body.clear()
                body.text = fc.get('back','')
            except Exception:
                pass

        bio = BytesIO()
        prs.save(bio)
        bio.seek(0)
        return bio.getvalue(), None
    except Exception as e:
        return None, f"PPTX export failed: {e}"


def export_study_pack_to_pptx(pack: dict, title_prefix="Study Pack"):
    """Export entire study pack to PPTX. Returns (bytes, error)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return None, "Missing dependency: install python-pptx (pip install python-pptx)"

    try:
        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = pack.get('session_name') or title_prefix

        # Topics slide
        topics = pack.get('topics') or []
        if topics:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = "Topics"
            tf = s.shapes.placeholders[1].text_frame
            tf.clear()
            for t in topics:
                p = tf.add_paragraph()
                p.text = str(t)
                p.level = 1

        # Summary slides (chunk into ~600 chars per slide)
        summary = (pack.get('summary') or "").strip()
        if summary:
            chunks = [summary[i:i+800] for i in range(0, len(summary), 800)]
            for idx, ch in enumerate(chunks, 1):
                layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                s = prs.slides.add_slide(layout)
                s.shapes.title.text = f"Summary (Part {idx})" if len(chunks) > 1 else "Summary"
                tf = s.shapes.placeholders[1].text_frame
                tf.clear()
                tf.text = ch

        # Flashcards as slides
        flashcards = pack.get('flashcards') or []
        if flashcards:
            for i, fc in enumerate(flashcards, 1):
                layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                s = prs.slides.add_slide(layout)
                s.shapes.title.text = f"Card {i}: {fc.get('front','')[:80]}"
                tf = s.shapes.placeholders[1].text_frame
                tf.clear()
                tf.text = fc.get('back','')

        # Quiz slides (question + choices)
        quiz = pack.get('quiz') or []
        if quiz:
            for i, q in enumerate(quiz, 1):
                layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                s = prs.slides.add_slide(layout)
                s.shapes.title.text = f"Q{i}. {q.get('question','')[:80]}"
                tf = s.shapes.placeholders[1].text_frame
                tf.clear()
                choices = q.get('choices') or []
                for cchoice in choices:
                    p = tf.add_paragraph()
                    p.text = f"- {cchoice}"
                    p.level = 1

        bio = BytesIO()
        prs.save(bio)
        bio.seek(0)
        return bio.getvalue(), None
    except Exception as e:
        return None, f"PPTX export failed: {e}"

# ============== TIER 1: AUDIO STUDY NOTES ==============
def generate_audio(text: str, voice_speed=1.0):
    """Convert text to audio using gTTS"""
    if not text:
        return None
    try:
        tts = gTTS(text=text, lang='en', slow=(voice_speed < 1.0))
        audio_bytes = BytesIO()
        tts.write_to_fp(audio_bytes)
        audio_bytes.seek(0)
        return audio_bytes
    except Exception as e:
        st.error(f"Audio generation failed: {e}")
        return None

# ============== TIER 2: SEMANTIC SEARCH (Q&A) ==============
def generate_embeddings(text: str):
    """Generate embeddings for semantic search"""
    if not client:
        return None
    try:
        response = client.embeddings.create(
            model="text-embedding-3-small",
            input=text
        )
        return response.data[0].embedding
    except Exception as e:
        st.error(f"Embedding generation failed: {e}")
        return None

def answer_question(context: str, question: str):
    """Answer a question based on document context"""
    if not context or not question:
        return ""
    prompt = f'''Answer this question based ONLY on the provided context. 
If the answer is not in the context, say "Not found in the provided text".
Include the relevant section from the context in your answer.

Context:
{truncate_chars(context, 3000)}

Question: {question}'''
    return openai_complete(prompt, temperature=0.3, max_tokens=500)

# ============== TIER 2: PRACTICE EXAM MODE ==============
def generate_timed_quiz(text: str, num_q=10, difficulty="medium"):
    """Generate timed quiz with difficulty levels"""
    if not text:
        return []
    
    difficulty_prompts = {
        "easy": "Create simple, basic questions",
        "medium": "Create moderate difficulty questions",
        "hard": "Create challenging, detailed questions"
    }
    
    prompt = f'''{difficulty_prompts.get(difficulty, "Create moderate difficulty questions")} from this text.
Create {num_q} multiple-choice questions.
Format as JSON array: [{{"question":"Q?","choices":["A","B","C","D"],"answer_index":0,"explanation":"..."}}]
Return ONLY JSON.

Text: {truncate_chars(text, 5000)}'''
    
    response = openai_complete(prompt, temperature=0.3, max_tokens=1500)
    try:
        if response:
            quiz = json.loads(response.replace("```json","").replace("```","").strip())
            return quiz
    except:
        st.error("Failed to generate exam")
    return []

# ============== TIER 3: GAMIFICATION ==============
def check_achievement(score_percentage):
    """Determine achievements based on score"""
    achievements = []
    if score_percentage == 100:
        achievements.append(("🏆 Perfect Score!", "Scored 100%!"))
    if score_percentage >= 80:
        achievements.append(("⭐ Star Student", "80% or higher"))
    if score_percentage >= 50:
        achievements.append(("📚 Learner", "Completed quiz"))
    return achievements

# ============== SESSION STATE ==============
init_db()

if "extracted_text" not in st.session_state:
    st.session_state.extracted_text = None
if "topics" not in st.session_state:
    st.session_state.topics = []
if "summary" not in st.session_state:
    st.session_state.summary = ""
if "quiz" not in st.session_state:
    st.session_state.quiz = []
if "flashcards" not in st.session_state:
    st.session_state.flashcards = []
if "streak" not in st.session_state:
    st.session_state.streak = 0
if "achievements" not in st.session_state:
    st.session_state.achievements = []
if "session_name" not in st.session_state:
    st.session_state.session_name = f"Study_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

# ============== SIDEBAR ==============
with st.sidebar:
    st.markdown("### ⚙️ Settings")
    # AI provider dropdown (OpenAI / Ollama / HuggingFace)
    if "provider" not in st.session_state:
        st.session_state.provider = "openai"
    provider = st.selectbox(
        "AI Provider:",
        ["openai", "ollama", "huggingface"],
        index=0,
        help="Choose LLM backend for summaries/quizzes"
    )
    st.session_state.provider = provider

    # model selection (used by OpenAI/OLLama)
    if "model" not in st.session_state:
        st.session_state.model = "gpt-4o-mini"
    model = st.selectbox("Model:", ["gpt-4o-mini", "gpt-4o"], index=0)
    st.session_state.model = model
    
    st.divider()
    
    st.markdown("### 📊 Progress Dashboard")
    history = get_quiz_history()
    if history:
        scores = [h[3]/h[4]*100 for h in history]
        avg_score = sum(scores) / len(scores) if scores else 0
        st.metric("Average Score", f"{avg_score:.1f}%")
        st.metric("Total Quizzes", len(history))
        
        if st.button("📈 View Progress Chart"):
            if PLOTLY_AVAILABLE:
                try:
                    fig = px.bar(x=list(range(len(scores))), y=scores, title="Quiz Score History")
                    st.plotly_chart(fig, use_container_width=True)
                except Exception as e:
                    # Plotly or its dependencies (pandas) may be broken/circular-importing.
                    # Fall back to matplotlib if available, otherwise show raw numbers.
                    st.warning(f"Plotly chart failed ({str(e)}). Showing fallback chart.")
                    try:
                        import matplotlib.pyplot as plt
                        fig2, ax = plt.subplots()
                        ax.bar(list(range(len(scores))), scores, color='tab:blue')
                        ax.set_title("Quiz Score History")
                        ax.set_xlabel("Attempt")
                        ax.set_ylabel("Score (%)")
                        st.pyplot(fig2)
                    except Exception:
                        st.write("Quiz scores:")
                        st.write(scores)
    
    st.divider()
    
    if st.session_state.achievements:
        st.markdown("### 🏆 Achievements Unlocked")
        for badge, desc in st.session_state.achievements:
            st.markdown(f'<div class="achievement">{badge}</div>', unsafe_allow_html=True)
    
    st.divider()
    st.caption("v2.0 Pro | Powered by OpenAI")

# ============== MAIN APP ==============
st.markdown("### 📄 Upload PDF")
uploaded = st.file_uploader("Choose PDF", type="pdf", label_visibility="collapsed")

if uploaded:
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown('<div class="status-badge done">✅ PDF</div>' if st.session_state.extracted_text else '<div class="status-badge pending">⏳ PDF</div>', unsafe_allow_html=True)
    with col2:
        st.markdown('<div class="status-badge done">✅ Topics</div>' if st.session_state.topics else '<div class="status-badge pending">⏳ Topics</div>', unsafe_allow_html=True)
    with col3:
        st.markdown('<div class="status-badge done">✅ Summary</div>' if st.session_state.summary else '<div class="status-badge pending">⏳ Summary</div>', unsafe_allow_html=True)
    with col4:
        st.markdown('<div class="status-badge done">✅ Resources</div>' if st.session_state.topics else '<div class="status-badge pending">⏳ Resources</div>', unsafe_allow_html=True)
    
    st.markdown("---")
    
    if st.button("⚡ Process PDF", use_container_width=True):
        if not client:
            st.error("🔑 OpenAI API Key Missing! Configure in Settings.")
        else:
            with st.spinner("📄 Extracting text..."):
                text = extract_pdf_text(uploaded)
                if text:
                    st.session_state.extracted_text = text
                    st.success("✅ Text extracted successfully")
                else:
                    st.error("❌ No text extracted")
            
            if st.session_state.extracted_text:
                with st.spinner("🎯 Detecting topics..."):
                    topics = ai_extract_topics(st.session_state.extracted_text)
                    st.session_state.topics = topics
                    if topics:
                        st.success(f"✅ Found {len(topics)} topics")
                
                with st.spinner("📝 Generating summary..."):
                    summary = ai_summarize(st.session_state.extracted_text)
                    st.session_state.summary = summary
                    if summary:
                        st.success("✅ Summary generated")
            
            st.rerun()
    
    st.markdown("---")
    
    # ============== ENHANCED TABS ==============
    tabs = st.tabs(["📝 Summary", "🎥 Videos", "🧪 Labs", "❓ Quiz", "🎴 Flashcards", "🎤 Audio", "❔ Q&A", "🎯 Exam", "📊 Progress", "📦 Export"])
    
    # TAB 0: Summary
    with tabs[0]:
        st.subheader("📝 Summary")
        if st.session_state.summary:
            st.markdown(st.session_state.summary)
            st.download_button("⬇️ Download MD", data=st.session_state.summary, file_name="summary.md", mime="text/markdown")
        else:
            st.info("Click **Process PDF** to generate summary")
    
    # TAB 1: Videos
    with tabs[1]:
        st.subheader("🎥 Video Resources")
        if st.session_state.topics:
            videos = fetch_videos_for_topics(st.session_state.topics)
            for topic, items in videos.items():
                st.markdown(f"**{topic}**")
                for v in items:
                    st.markdown(f"[{v['title']}]({v['url']})")
        else:
            st.info("Process PDF to detect topics")
    
    # TAB 2: Interactive Labs
    with tabs[2]:
        st.subheader("🧪 Hands-On Labs & Projects")
        if st.session_state.topics and st.session_state.extracted_text:
            # Add caching for lab content to avoid regenerating
            if "lab_content" not in st.session_state:
                st.session_state.lab_content = {}
            
            for topic in st.session_state.topics:
                st.markdown(f"### {topic}")
                
                # Check if lab content is cached
                if topic not in st.session_state.lab_content:
                    if st.button(f"📚 Generate Lab Content for {topic}", key=f"lab_{topic}"):
                        with st.spinner(f"Generating lab content for {topic}..."):
                            lab_data = generate_lab_content(st.session_state.extracted_text, topic)
                            st.session_state.lab_content[topic] = lab_data
                            st.rerun()
                
                if topic in st.session_state.lab_content:
                    lab_data = st.session_state.lab_content[topic]
                    
                    # Real-world applications
                    if lab_data.get("real_world"):
                        st.markdown("#### 🌍 Real-World Applications")
                        st.markdown(lab_data["real_world"])
                    
                    # Project ideas
                    projects = lab_data.get("projects", [])
                    if projects:
                        st.markdown("#### 🎯 Project Ideas")
                        for i, proj in enumerate(projects, 1):
                            st.markdown(f"**Project {i}:** {proj}")
                    
                    # Hands-on assignments
                    assignments = lab_data.get("assignments", [])
                    if assignments:
                        st.markdown("#### ✍️ Hands-On Assignments")
                        for i, assign in enumerate(assignments, 1):
                            st.markdown(f"**Assignment {i}:** {assign}")
                    
                    # Lab tips
                    tips = lab_data.get("lab_tips", [])
                    if tips:
                        st.markdown("#### 💡 Lab Tips")
                        for i, tip in enumerate(tips, 1):
                            st.markdown(f"**Tip {i}:** {tip}")
                
                st.markdown("---")
        else:
            st.info("Process PDF to detect topics and extract content")
    
    # TAB 3: Regular Quiz
    with tabs[3]:
        st.subheader("❓ Practice Quiz")
        if st.session_state.extracted_text:
            num_q = st.slider("Questions:", 3, 15, 5)
            if st.button("🎲 Generate Quiz"):
                with st.spinner("Generating quiz..."):
                    prompt = f'''Create {num_q} multiple-choice questions from this text.
Format as JSON array: [{{"question":"Q?","choices":["A","B","C","D"],"answer_index":0,"explanation":"..."}}]
Return ONLY JSON.

Text: {truncate_chars(st.session_state.extracted_text, 5000)}'''
                    response = openai_complete(prompt, temperature=0.2, max_tokens=1500)
                    try:
                        quiz = json.loads(response.replace("```json","").replace("```","").strip())
                        st.session_state.quiz = quiz
                    except:
                        st.error("Failed to generate quiz")
            
            if st.session_state.quiz:
                score = 0
                total = len(st.session_state.quiz)
                for i, q in enumerate(st.session_state.quiz, 1):
                    st.markdown(f"**Q{i}. {q['question']}**")
                    ans = st.radio(f"Q{i}", q["choices"], key=f"q_{i}", label_visibility="collapsed")
                    if st.button("Check", key=f"btn_{i}"):
                        correct = q["choices"][q["answer_index"]]
                        if ans == correct:
                            st.success("✅ Correct!")
                            score += 1
                        else:
                            st.error(f"❌ Correct answer: {correct}")
                        st.info(f"**Explanation:** {q.get('explanation', 'N/A')}")
        else:
            st.info("Process PDF first")
    
    # TAB 4: TIER 1 - Flashcards
    with tabs[4]:
        st.subheader("🎴 Flashcard Generator")
        if st.session_state.extracted_text:
            num_cards = st.slider("Number of flashcards:", 5, 30, 10)
            if st.button("📚 Generate Flashcards"):
                with st.spinner("Creating flashcards..."):
                    cards = generate_flashcards(st.session_state.extracted_text, num_cards)
                    if cards:
                        st.session_state.flashcards = cards
                        st.success(f"✅ Created {len(cards)} flashcards")
            
            if st.session_state.flashcards:
                st.markdown("### Study Your Flashcards")
                if "card_index" not in st.session_state:
                    st.session_state.card_index = 0
                if "card_flipped" not in st.session_state:
                    st.session_state.card_flipped = False
                
                card = st.session_state.flashcards[st.session_state.card_index]
                
                col1, col2, col3 = st.columns([1, 2, 1])
                with col1:
                    if st.button("⬅️ Previous"):
                        st.session_state.card_index = max(0, st.session_state.card_index - 1)
                        st.session_state.card_flipped = False
                        st.rerun()
                with col2:
                    st.markdown(f"<div class='flashcard'>{card['front'] if not st.session_state.card_flipped else card['back']}</div>", unsafe_allow_html=True)
                with col3:
                    if st.button("Next ➡️"):
                        st.session_state.card_index = min(len(st.session_state.flashcards) - 1, st.session_state.card_index + 1)
                        st.session_state.card_flipped = False
                        st.rerun()
                
                st.button("🔄 Flip Card", key="flip", on_click=lambda: st.session_state.update({"card_flipped": not st.session_state.card_flipped}))
                st.caption(f"Card {st.session_state.card_index + 1} of {len(st.session_state.flashcards)}")
                
                # Export flashcards
                anki_format = export_anki_format(st.session_state.flashcards)
                st.download_button("📥 Download (Anki Format)", data=anki_format, file_name="flashcards.txt", mime="text/plain")
                # DOCX export
                docx_bytes, docx_err = export_to_docx(st.session_state.flashcards)
                if docx_bytes:
                    st.download_button("📥 Download (Word .docx)", data=docx_bytes, file_name="flashcards.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                else:
                    if docx_err:
                        st.caption(docx_err)

                # PDF export
                pdf_bytes, pdf_err = export_to_pdf(st.session_state.flashcards)
                if pdf_bytes:
                    st.download_button("📥 Download (PDF)", data=pdf_bytes, file_name="flashcards.pdf", mime="application/pdf")
                else:
                    if pdf_err:
                        st.caption(pdf_err)
                # PPTX export for flashcards
                pptx_bytes, pptx_err = export_flashcards_to_pptx(st.session_state.flashcards)
                if pptx_bytes:
                    st.download_button("📥 Download (PPTX)", data=pptx_bytes, file_name="flashcards.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                else:
                    if pptx_err:
                        st.caption(pptx_err)
        else:
            st.info("Process PDF first")
    
    # TAB 5: TIER 1 - Audio Notes
    with tabs[5]:
        st.subheader("🎤 Audio Study Notes")
        if st.session_state.summary:
            st.markdown("Convert your summary to audio for listening on-the-go!")
            voice_speed = st.slider("Speed:", 0.5, 2.0, 1.0)
            
            if st.button("🎙️ Generate Audio"):
                with st.spinner("Generating audio..."):
                    audio = generate_audio(st.session_state.summary, voice_speed)
                    if audio:
                        st.audio(audio, format="audio/mp3")
                        st.download_button("⬇️ Download MP3", data=audio, file_name="study_notes.mp3", mime="audio/mp3")
        else:
            st.info("Generate summary first")
    
    # TAB 6: TIER 2 - Q&A (Semantic Search)
    with tabs[6]:
        st.subheader("❔ Ask Questions About Your Notes")
        if st.session_state.extracted_text:
            question = st.text_input("💭 Ask a question about the content...")
            if question and st.button("🔍 Search"):
                with st.spinner("Finding answer..."):
                    answer = answer_question(st.session_state.extracted_text, question)
                    if answer:
                        st.markdown(f"**Answer:**\n{answer}")
        else:
            st.info("Process PDF first")
    
    # TAB 7: TIER 2 - Practice Exam
    with tabs[7]:
        st.subheader("🎯 Practice Exam Mode")
        if st.session_state.extracted_text:
            col1, col2, col3 = st.columns(3)
            with col1:
                num_exam_q = st.slider("Questions:", 5, 20, 10)
            with col2:
                difficulty = st.selectbox("Difficulty:", ["easy", "medium", "hard"])
            with col3:
                time_limit = st.slider("Time (min):", 5, 60, 15)
            
            if st.button("🚀 Start Exam"):
                with st.spinner("Generating exam..."):
                    exam_quiz = generate_timed_quiz(st.session_state.extracted_text, num_exam_q, difficulty)
                    if exam_quiz:
                        st.session_state.exam_quiz = exam_quiz
                        st.session_state.exam_start_time = time.time()
                        st.session_state.time_limit = time_limit * 60
                        st.success("Exam started! Answer all questions.")
            
            if "exam_quiz" in st.session_state:
                # Create a placeholder for timer that updates dynamically
                timer_placeholder = st.empty()
                
                elapsed = time.time() - st.session_state.exam_start_time
                remaining = max(0, st.session_state.time_limit - elapsed)
                
                # Display timer with color coding
                if remaining > 300:  # More than 5 minutes
                    timer_placeholder.info(f"⏱️ Time remaining: {int(remaining//60)}:{int(remaining%60):02d}")
                elif remaining > 60:  # Between 1-5 minutes
                    timer_placeholder.warning(f"⏱️ Time remaining: {int(remaining//60)}:{int(remaining%60):02d}")
                else:  # Less than 1 minute
                    timer_placeholder.error(f"⏱️ Time remaining: {int(remaining//60)}:{int(remaining%60):02d}")
                
                if remaining <= 0:
                    st.error("⏰ Time's up! Exam has ended.")
                    st.stop()
                else:
                    score = 0
                    for i, q in enumerate(st.session_state.exam_quiz, 1):
                        st.markdown(f"**Q{i}. {q['question']}**")
                        ans = st.radio(f"Q{i}", q["choices"], key=f"exam_q_{i}", label_visibility="collapsed")
                        correct = q["choices"][q["answer_index"]]
                        if ans == correct:
                            score += 1
                    
                    if st.button("📊 Submit Exam"):
                        percentage = (score / len(st.session_state.exam_quiz)) * 100
                        st.success(f"✅ Score: {score}/{len(st.session_state.exam_quiz)} ({percentage:.1f}%)")
                        
                        # Check achievements
                        achievements = check_achievement(percentage)
                        for badge, desc in achievements:
                            if badge not in [a[0] for a in st.session_state.achievements]:
                                st.session_state.achievements.append((badge, desc))
                                st.balloons()
                        
                        add_quiz_score(st.session_state.session_name, "Exam", score, len(st.session_state.exam_quiz))
        else:
            st.info("Process PDF first")
    
    # TAB 8: TIER 3 - Progress Dashboard
    with tabs[8]:
        st.subheader("📊 Your Learning Progress")
        history = get_quiz_history()
        if history:
            scores = [h[3]/h[4]*100 for h in history]
            
            # Show metrics
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Total Quizzes", len(history))
            with col2:
                st.metric("Average Score", f"{sum(scores)/len(scores):.1f}%")
            with col3:
                st.metric("Best Score", f"{max(scores):.1f}%")
            with col4:
                st.metric("Streak", st.session_state.streak)
            
            # Show chart
            if PLOTLY_AVAILABLE:
                fig = go.Figure()
                fig.add_trace(go.Scatter(y=scores, mode='lines+markers', name='Score', line=dict(color='#6366f1')))
                fig.update_layout(title="Quiz Score Progress", xaxis_title="Quiz #", yaxis_title="Score %")
                st.plotly_chart(fig, use_container_width=True)
            else:
                st.bar_chart(scores)
            
            # Show achievements
            if st.session_state.achievements:
                st.markdown("### 🏆 Achievements")
                for badge, desc in st.session_state.achievements:
                    st.markdown(f'<div class="achievement">{badge}</div>', unsafe_allow_html=True)
        else:
            st.info("No quiz history yet. Take some quizzes to see progress!")
    
    # TAB 9: Export
    with tabs[9]:
        st.subheader("📦 Export Study Pack")
        
        export_data = {
            "session_name": st.session_state.session_name,
            "topics": st.session_state.topics,
            "summary": st.session_state.summary,
            "flashcards": st.session_state.flashcards,
            "quiz": st.session_state.quiz,
            "created_at": datetime.now().isoformat(),
        }
        
        col1, col2 = st.columns(2)
        with col1:
            # Offer DOCX and PDF exports for the study pack (instead of raw JSON)
            docx_bytes, docx_err = export_study_pack_to_docx(export_data)
            if docx_bytes:
                st.download_button("📥 Download Study Pack (Word .docx)", data=docx_bytes, file_name="study_pack.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            else:
                if docx_err:
                    st.caption(docx_err)

            pdf_bytes, pdf_err = export_study_pack_to_pdf(export_data)
            if pdf_bytes:
                st.download_button("📥 Download Study Pack (PDF)", data=pdf_bytes, file_name="study_pack.pdf", mime="application/pdf")
            else:
                if pdf_err:
                    st.caption(pdf_err)
            # PPTX export for full study pack
            pptx_bytes, pptx_err = export_study_pack_to_pptx(export_data)
            if pptx_bytes:
                st.download_button("📥 Download Study Pack (PPTX)", data=pptx_bytes, file_name="study_pack.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            else:
                if pptx_err:
                    st.caption(pptx_err)
        with col2:
            if st.session_state.flashcards:
                anki = export_anki_format(st.session_state.flashcards)
                st.download_button(
                    "📚 Anki Flashcards",
                    data=anki,
                    file_name="flashcards.txt"
                )

else:
    st.info("👆 Upload a PDF to start learning!")

st.markdown("---")
st.caption("🚀 AI Mentor v2.0 Pro | All Tier Features Included | Powered by OpenAI")
